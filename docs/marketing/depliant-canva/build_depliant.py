#!/usr/bin/env python3
"""Génère le dépliant et la présentation Somafrik importables dans Canva.

Les captures sont des screenshots runtime officiels (données fictives).
Le PPTX conserve textes et formes éditables : Fichier → Importer dans Canva.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
PREVIEWS = ROOT / "previews"
DIST = ROOT

# Identité officielle (vitrine + design tokens)
INK = RGBColor(0x0F, 0x17, 0x2A)
BRAND = RGBColor(0x1D, 0x4E, 0xD8)
BRAND_50 = RGBColor(0xEF, 0xF6, 0xFF)
BRAND_100 = RGBColor(0xDB, 0xEA, 0xFE)
BRAND_700 = RGBColor(0x1E, 0x40, 0xAF)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
CANVAS = RGBColor(0xF7, 0xF9, 0xFC)
MUTED = RGBColor(0x64, 0x73, 0x8B)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xB4, 0x53, 0x09)
NAVY = RGBColor(0x0B, 0x12, 0x20)

PIL_INK = (15, 23, 42)
PIL_BRAND = (29, 78, 216)
PIL_BRAND_50 = (239, 246, 255)
PIL_BRAND_100 = (219, 234, 254)
PIL_BRAND_700 = (30, 64, 175)
PIL_TEAL = (15, 118, 110)
PIL_CANVAS = (247, 249, 252)
PIL_MUTED = (100, 115, 139)
PIL_LINE = (226, 232, 240)
PIL_WHITE = (255, 255, 255)
PIL_NAVY = (11, 18, 32)

FONT_REG = "/usr/share/fonts/truetype/macos/Inter-Regular.ttf"
FONT_MED = "/usr/share/fonts/truetype/macos/Inter-Medium.ttf"
FONT_SEMI = "/usr/share/fonts/truetype/macos/Inter-SemiBold.ttf"
FONT_BOLD = "/usr/share/fonts/truetype/macos/Inter-Bold.ttf"

WEB = {
    "login": ASSETS / "web" / "01-connexion-etablissement.png",
    "dashboard": ASSETS / "web" / "02-tableau-de-bord-etablissement.png",
    "classes": ASSETS / "web" / "03-classes-liste.png",
    "eleves": ASSETS / "web" / "05-eleves-annuaire.png",
    "dossier": ASSETS / "web" / "06-eleve-dossier.png",
}
MOBILE = {
    "classes": ASSETS / "mobile" / "02-classes-liste.png",
    "eleves": ASSETS / "mobile" / "04-eleves-liste.png",
    "enseignants": ASSETS / "mobile" / "07-enseignants.png",
    "paiements": ASSETS / "mobile" / "12-paiements.png",
    "appel": ASSETS / "mobile" / "16-presences-appel.png",
    "notes": ASSETS / "mobile" / "19-notes-saisie.png",
    "parent": ASSETS / "mobile" / "20-parent-accueil.png",
}
LOGO = ASSETS / "somafrik-logo.png"
ICON = ASSETS / "somafrik-icon.png"

FEATURES = [
    ("Scolarité", "Élèves, classes et enseignants au même endroit."),
    ("Pédagogie", "Présences, notes et évaluations du quotidien."),
    ("Finances", "Frais scolaires, paiements et suivi des soldes."),
    ("Communication", "Messages et annonces de l’établissement."),
    ("Pilotage", "Tableaux de bord et indicateurs par école."),
]
AUDIENCES = [
    ("Direction", "Suivre classes, effectifs, présences et frais, puis décider."),
    ("Administration", "Inscrire les élèves, organiser les classes, enregistrer les paiements."),
    ("Enseignants", "Faire l’appel et saisir les notes, sur web ou mobile."),
    ("Parents", "Consulter notes, présences et frais, selon les accès accordés."),
]
BENEFITS = [
    "Centraliser dossiers, classes et comptes",
    "Suivre la scolarité sans tableurs dispersés",
    "Piloter les frais et le reste à payer",
    "Appel et notes sur le terrain, depuis le mobile",
    "Les familles restent informées",
]


def font(path: str, size: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(path, size)


def rounded_mask(size: tuple[int, int], radius: int) -> Image.Image:
    mask = Image.new("L", size, 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, size[0] - 1, size[1] - 1), radius=radius, fill=255)
    return mask


def frame_shot(src: Path, dest: Path, *, radius: int = 28, max_w: int = 1600) -> Path:
    dest.parent.mkdir(parents=True, exist_ok=True)
    im = Image.open(src).convert("RGBA")
    if im.width > max_w:
        h = int(im.height * max_w / im.width)
        im = im.resize((max_w, h), Image.Resampling.LANCZOS)
    r = min(radius, im.width // 8, im.height // 8)
    mask = rounded_mask(im.size, r)
    im.putalpha(mask)
    pad = 18
    canvas = Image.new("RGBA", (im.width + pad * 2, im.height + pad * 2), (0, 0, 0, 0))
    shadow = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    sd = ImageDraw.Draw(shadow)
    sd.rounded_rectangle((pad, pad + 6, pad + im.width, pad + im.height + 8), radius=r, fill=(15, 23, 42, 55))
    shadow = shadow.filter(ImageFilter.GaussianBlur(10))
    canvas = Image.alpha_composite(canvas, shadow)
    canvas.paste(im, (pad, pad), im)
    canvas.save(dest, "PNG")
    return dest


def framed(key: str, src: Path, max_w: int = 1600, radius: int = 28) -> Path:
    dest = Path("/tmp/somafrik-depliant") / f"{key}.png"
    if dest.exists() and dest.stat().st_mtime >= src.stat().st_mtime:
        return dest
    return frame_shot(src, dest, radius=radius, max_w=max_w)


# ---------------------------------------------------------------------------
# PowerPoint helpers (éléments natifs = éditables dans Canva)
# ---------------------------------------------------------------------------

def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color: RGBColor, *, rounded: bool = False, radius: float = 0.08):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, x, y, w, h)
    _set_fill(shape, color)
    if rounded:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text: str,
    *,
    size: float = 14,
    bold: bool = False,
    color: RGBColor = INK,
    align=PP_ALIGN.LEFT,
    font_name: str = "Inter",
    anchor: str = "t",
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    # anchor: t / m / b
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    body = tf._txBody
    bodyPr = body.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", {"t": "t", "m": "ctr", "b": "b"}[anchor])
    return box


def add_lines(
    slide,
    x,
    y,
    w,
    h,
    lines: list[tuple[str, dict]],
    *,
    align=PP_ALIGN.LEFT,
    spacing: float = 1.0,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, style) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = style.get("align", align)
        p.space_after = Pt(style.get("after", 4))
        p.line_spacing = spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(style.get("size", 14))
        run.font.bold = style.get("bold", False)
        run.font.color.rgb = style.get("color", INK)
        run.font.name = style.get("font", "Inter")
    return box


def add_pic(slide, path: Path, x, y, w=None, h=None):
    kwargs = {}
    if w is not None:
        kwargs["width"] = w
    if h is not None:
        kwargs["height"] = h
    return slide.shapes.add_picture(str(path), x, y, **kwargs)


def blank(prs) -> object:
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# Présentation 16:9 — 6 diapos éditables
# ---------------------------------------------------------------------------

def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height

    # --- 1. Couverture ---
    s = blank(prs)
    add_rect(s, 0, 0, Inches(5.7), H, NAVY)
    add_rect(s, Inches(5.7), 0, W - Inches(5.7), H, CANVAS)
    add_rect(s, Inches(5.55), 0, Inches(0.22), H, BRAND)

    add_pic(s, ICON, Inches(0.45), Inches(0.38), Inches(0.72), Inches(0.72))
    add_text(s, Inches(1.28), Inches(0.42), Inches(4.0), Inches(0.32), "SOMAFRIK", size=13, bold=True, color=WHITE)
    add_text(
        s,
        Inches(1.28),
        Inches(0.70),
        Inches(4.0),
        Inches(0.28),
        "Transformer l’éducation, renforcer l’impact",
        size=10,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    add_text(
        s,
        Inches(0.45),
        Inches(1.55),
        Inches(4.8),
        Inches(0.32),
        "GESTION SCOLAIRE MODERNE",
        size=12,
        bold=True,
        color=RGBColor(0x93, 0xC5, 0xFD),
    )
    add_text(
        s,
        Inches(0.45),
        Inches(1.95),
        Inches(4.9),
        Inches(2.4),
        "Pilotez votre établissement depuis un seul endroit",
        size=32,
        bold=True,
        color=WHITE,
    )
    add_text(
        s,
        Inches(0.45),
        Inches(4.45),
        Inches(4.85),
        Inches(1.35),
        "Somafrik réunit les élèves, les enseignants, les présences, les notes, les finances et la communication dans une plateforme accessible sur Web et Mobile.",
        size=14,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    add_rect(s, Inches(0.45), Inches(6.05), Inches(2.55), Inches(0.52), BRAND, rounded=True, radius=0.2)
    add_text(
        s,
        Inches(0.45),
        Inches(6.08),
        Inches(2.55),
        Inches(0.48),
        "somafrik.app",
        size=15,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        anchor="m",
    )
    add_text(
        s,
        Inches(3.15),
        Inches(6.15),
        Inches(2.2),
        Inches(0.4),
        "Web  ·  Mobile",
        size=13,
        bold=True,
        color=RGBColor(0x93, 0xC5, 0xFD),
        anchor="m",
    )

    dash = framed("web-dashboard", WEB["dashboard"], max_w=1400)
    add_pic(s, dash, Inches(6.05), Inches(0.85), Inches(6.85))
    add_text(
        s,
        Inches(6.15),
        Inches(6.85),
        Inches(6.7),
        Inches(0.4),
        "Application web — tableau de bord établissement (données fictives)",
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    # --- 2. Application web ---
    s = blank(prs)
    add_rect(s, 0, 0, W, H, WHITE)
    add_rect(s, 0, 0, W, Inches(0.08), BRAND)
    add_text(s, Inches(0.5), Inches(0.28), Inches(4), Inches(0.28), "APPLICATION WEB", size=12, bold=True, color=BRAND)
    add_text(
        s,
        Inches(0.5),
        Inches(0.55),
        Inches(12.3),
        Inches(0.55),
        "Ce que voit la direction une fois connectée",
        size=26,
        bold=True,
        color=INK,
    )
    add_text(
        s,
        Inches(0.5),
        Inches(1.12),
        Inches(12.3),
        Inches(0.4),
        "Annuaire des élèves, organisation des classes, suivi des enseignants et tableaux de bord — depuis un écran complet.",
        size=13,
        color=MUTED,
    )

    web_items = [
        (WEB["dashboard"], "Tableau de bord"),
        (WEB["classes"], "Classes"),
        (WEB["eleves"], "Annuaire des élèves"),
    ]
    left = Inches(0.45)
    gap = Inches(0.22)
    card_w = Inches(4.05)
    for i, (src, caption) in enumerate(web_items):
        x = left + i * (card_w + gap)
        add_rect(s, x, Inches(1.65), card_w, Inches(5.35), CANVAS, rounded=True, radius=0.06)
        shot = framed(f"web-{i}", src, max_w=1100, radius=18)
        add_pic(s, shot, x + Inches(0.12), Inches(1.78), card_w - Inches(0.24))
        add_text(
            s,
            x,
            Inches(6.55),
            card_w,
            Inches(0.32),
            caption,
            size=13,
            bold=True,
            color=INK,
            align=PP_ALIGN.CENTER,
        )

    # --- 3. Application mobile ---
    s = blank(prs)
    add_rect(s, 0, 0, W, H, CANVAS)
    add_rect(s, 0, 0, W, Inches(0.08), TEAL)
    add_text(s, Inches(0.5), Inches(0.28), Inches(6), Inches(0.28), "APPLICATION MOBILE", size=12, bold=True, color=TEAL)
    add_text(
        s,
        Inches(0.5),
        Inches(0.55),
        Inches(12.3),
        Inches(0.5),
        "Le bureau pour piloter, le mobile pour le terrain",
        size=26,
        bold=True,
        color=INK,
    )
    add_text(
        s,
        Inches(0.5),
        Inches(1.1),
        Inches(12.3),
        Inches(0.35),
        "Les enseignants font l’appel et saisissent les notes. Les parents suivent la scolarité de l’enfant.",
        size=13,
        color=MUTED,
    )

    phones = [
        (MOBILE["classes"], "Classes"),
        (MOBILE["appel"], "Appel des présences"),
        (MOBILE["notes"], "Saisie des notes"),
        (MOBILE["parent"], "Espace parent"),
    ]
    left = Inches(0.55)
    gap = Inches(0.28)
    phone_w = Inches(2.85)
    for i, (src, caption) in enumerate(phones):
        x = left + i * (phone_w + gap)
        add_rect(s, x, Inches(1.58), phone_w, Inches(5.45), WHITE, rounded=True, radius=0.1)
        shot = framed(f"m-{i}", src, max_w=720, radius=36)
        add_pic(s, shot, x + Inches(0.12), Inches(1.7), phone_w - Inches(0.24))
        add_text(
            s,
            x,
            Inches(6.95),
            phone_w,
            Inches(0.32),
            caption,
            size=12,
            bold=True,
            color=INK,
            align=PP_ALIGN.CENTER,
        )

    # --- 4. Preuves métier ---
    s = blank(prs)
    add_rect(s, 0, 0, W, H, WHITE)
    add_rect(s, 0, 0, W, Inches(0.08), BRAND)
    add_text(s, Inches(0.5), Inches(0.28), Inches(6), Inches(0.28), "PREUVES MÉTIER", size=12, bold=True, color=BRAND)
    add_text(
        s,
        Inches(0.5),
        Inches(0.55),
        Inches(12.3),
        Inches(0.5),
        "Somafrik dans le quotidien de l’établissement",
        size=26,
        bold=True,
        color=INK,
    )

    proofs = [
        (WEB["login"], "Connexion web", "Code établissement et identifiant."),
        (WEB["dossier"], "Dossier élève", "Identité, inscription, responsables."),
        (MOBILE["paiements"], "Paiements", "Frais, encaissements, reçus."),
        (MOBILE["eleves"], "Élèves (mobile)", "Annuaire et suivi de classe."),
    ]
    left = Inches(0.4)
    gap = Inches(0.2)
    card_w = Inches(3.08)
    for i, (src, title, desc) in enumerate(proofs):
        x = left + i * (card_w + gap)
        add_rect(s, x, Inches(1.25), card_w, Inches(5.85), CANVAS, rounded=True, radius=0.07)
        shot = framed(f"p-{i}", src, max_w=900, radius=20)
        add_pic(s, shot, x + Inches(0.12), Inches(1.4), card_w - Inches(0.24))
        add_text(s, x + Inches(0.16), Inches(6.35), card_w - Inches(0.3), Inches(0.3), title, size=13, bold=True, color=INK)
        add_text(s, x + Inches(0.16), Inches(6.65), card_w - Inches(0.3), Inches(0.32), desc, size=11, color=MUTED)

    # --- 5. Fonctionnalités + publics ---
    s = blank(prs)
    add_rect(s, 0, 0, W, H, CANVAS)
    add_rect(s, 0, 0, W, Inches(0.08), BRAND)
    add_text(s, Inches(0.5), Inches(0.25), Inches(8), Inches(0.26), "FONCTIONNALITÉS", size=12, bold=True, color=BRAND)
    add_text(
        s,
        Inches(0.5),
        Inches(0.5),
        Inches(12.3),
        Inches(0.45),
        "Les besoins couverts aujourd’hui",
        size=26,
        bold=True,
        color=INK,
    )

    card_w = Inches(2.38)
    gap = Inches(0.16)
    left = Inches(0.45)
    for i, (title, desc) in enumerate(FEATURES):
        x = left + i * (card_w + gap)
        add_rect(s, x, Inches(1.15), card_w, Inches(2.05), WHITE, rounded=True, radius=0.1)
        add_rect(s, x, Inches(1.15), Inches(0.1), Inches(2.05), BRAND)
        add_text(s, x + Inches(0.22), Inches(1.3), card_w - Inches(0.35), Inches(0.4), title, size=15, bold=True, color=INK)
        add_text(s, x + Inches(0.22), Inches(1.75), card_w - Inches(0.35), Inches(1.2), desc, size=12, color=MUTED)

    add_text(s, Inches(0.5), Inches(3.45), Inches(8), Inches(0.28), "POUR QUI ?", size=12, bold=True, color=BRAND)
    add_text(
        s,
        Inches(0.5),
        Inches(3.7),
        Inches(12),
        Inches(0.4),
        "Pensé d’abord pour l’établissement — chaque profil retrouve son travail quotidien.",
        size=16,
        bold=True,
        color=INK,
    )
    card_w = Inches(3.05)
    gap = Inches(0.18)
    left = Inches(0.45)
    for i, (title, desc) in enumerate(AUDIENCES):
        x = left + i * (card_w + gap)
        add_rect(s, x, Inches(4.25), card_w, Inches(2.75), WHITE, rounded=True, radius=0.08)
        add_text(s, x + Inches(0.22), Inches(4.45), card_w - Inches(0.4), Inches(0.4), title, size=16, bold=True, color=BRAND)
        add_text(s, x + Inches(0.22), Inches(4.95), card_w - Inches(0.4), Inches(1.7), desc, size=13, color=INK)

    # --- 6. Clôture ---
    s = blank(prs)
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, 0, Inches(0.18), H, BRAND)
    add_pic(s, ICON, Inches(0.7), Inches(0.55), Inches(0.7), Inches(0.7))
    add_text(s, Inches(1.55), Inches(0.62), Inches(6), Inches(0.28), "SOMAFRIK", size=14, bold=True, color=WHITE)
    add_text(
        s,
        Inches(1.55),
        Inches(0.92),
        Inches(8),
        Inches(0.28),
        "La plateforme qui simplifie la gestion de votre établissement scolaire.",
        size=12,
        color=RGBColor(0x94, 0xA3, 0xB8),
    )
    add_text(
        s,
        Inches(0.7),
        Inches(1.7),
        Inches(7.4),
        Inches(1.6),
        "Ouvrez l’espace de votre établissement",
        size=32,
        bold=True,
        color=WHITE,
    )
    add_text(
        s,
        Inches(0.7),
        Inches(3.4),
        Inches(7.2),
        Inches(1.2),
        "Si votre école dispose déjà d’un compte Somafrik, connectez-vous pour retrouver les classes, les élèves, les présences, les notes et les finances.",
        size=15,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    add_rect(s, Inches(0.7), Inches(4.85), Inches(3.15), Inches(0.58), BRAND, rounded=True, radius=0.18)
    add_text(
        s,
        Inches(0.7),
        Inches(4.88),
        Inches(3.15),
        Inches(0.52),
        "somafrik.app /connexion",
        size=14,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        anchor="m",
    )
    add_text(
        s,
        Inches(4.05),
        Inches(4.95),
        Inches(4.2),
        Inches(0.4),
        "somafrik.app",
        size=16,
        bold=True,
        color=RGBColor(0x93, 0xC5, 0xFD),
        anchor="m",
    )

    add_rect(s, Inches(8.55), Inches(0.55), Inches(4.25), Inches(6.4), RGBColor(0x12, 0x1A, 0x2C), rounded=True, radius=0.06)
    add_text(s, Inches(8.8), Inches(0.75), Inches(3.85), Inches(0.3), "Chaque établissement reste dans son périmètre", size=13, bold=True, color=WHITE)
    security = [
        ("Authentification", "Connexion avec le code de l’établissement et un identifiant."),
        ("Séparation", "Les données d’une école ne sont pas mélangées avec celles d’une autre."),
        ("Rôles", "Direction, administration, enseignant ou parent : chacun voit son périmètre."),
    ]
    for i, (t, d) in enumerate(security):
        y = Inches(1.3) + i * Inches(1.35)
        add_text(s, Inches(8.8), y, Inches(3.85), Inches(0.32), t, size=14, bold=True, color=RGBColor(0x93, 0xC5, 0xFD))
        add_text(s, Inches(8.8), y + Inches(0.35), Inches(3.85), Inches(0.85), d, size=12, color=RGBColor(0xCB, 0xD5, 0xE1))

    add_text(
        s,
        Inches(0.7),
        Inches(6.85),
        Inches(12),
        Inches(0.3),
        "© 2026 Somafrik  ·  Captures runtime, données fictives  ·  Document modifiable dans Canva",
        size=11,
        color=RGBColor(0x64, 0x74, 0x8B),
    )

    out = DIST / "Somafrik-presentation-Canva.pptx"
    prs.save(out)
    return out


# ---------------------------------------------------------------------------
# Dépliant 3 volets A4 paysage (recto / verso)
# ---------------------------------------------------------------------------

def build_leaflet() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(11.693)
    prs.slide_height = Inches(8.268)
    W, H = prs.slide_width, prs.slide_height
    panel = W / 3

    def fold_guides(slide):
        for i in (1, 2):
            add_rect(slide, int(panel * i) - Emu(20000), 0, Emu(40000), H, RGBColor(0xCB, 0xD5, 0xE1))

    # RECTO — extérieur : dos | couverture | rabat
    s = blank(prs)
    add_rect(s, 0, 0, panel, H, NAVY)
    add_rect(s, panel, 0, panel, H, BRAND)
    add_rect(s, panel * 2, 0, panel, H, CANVAS)

    # Dos (gauche)
    add_pic(s, ICON, Inches(0.28), Inches(0.32), Inches(0.58), Inches(0.58))
    add_text(s, Inches(0.95), Inches(0.38), Inches(2.6), Inches(0.28), "SOMAFRIK", size=13, bold=True, color=WHITE)
    add_text(s, Inches(0.95), Inches(0.64), Inches(2.6), Inches(0.28), "Web et mobile", size=10, color=RGBColor(0x93, 0xC5, 0xFD))
    add_text(s, Inches(0.28), Inches(1.2), Inches(3.35), Inches(0.7), "Contact", size=18, bold=True, color=WHITE)
    add_text(s, Inches(0.28), Inches(1.9), Inches(3.35), Inches(0.7), "somafrik.app\nsomafrik.app/connexion", size=13, color=RGBColor(0xE2, 0xE8, 0xF0))
    add_text(s, Inches(0.28), Inches(2.85), Inches(3.35), Inches(0.35), "Sécurité", size=14, bold=True, color=RGBColor(0x93, 0xC5, 0xFD))
    add_text(
        s,
        Inches(0.28),
        Inches(3.25),
        Inches(3.35),
        Inches(2.6),
        "Authentification par établissement.\n\nDonnées séparées par école.\n\nAccès selon le rôle : direction, administration, enseignant, parent.",
        size=12,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    add_text(
        s,
        Inches(0.28),
        Inches(7.45),
        Inches(3.35),
        Inches(0.5),
        "© 2026 Somafrik\nCaptures runtime · données fictives",
        size=9,
        color=RGBColor(0x64, 0x74, 0x8B),
    )

    # Couverture (centre)
    add_pic(s, ICON, panel + Inches(0.28), Inches(0.35), Inches(0.62), Inches(0.62))
    add_text(
        s,
        panel + Inches(1.0),
        Inches(0.42),
        Inches(2.6),
        Inches(0.28),
        "SOMAFRIK",
        size=13,
        bold=True,
        color=WHITE,
    )
    add_text(
        s,
        panel + Inches(1.0),
        Inches(0.68),
        Inches(2.6),
        Inches(0.28),
        "Gestion scolaire moderne",
        size=10,
        color=RGBColor(0xDB, 0xEA, 0xFE),
    )
    add_text(
        s,
        panel + Inches(0.28),
        Inches(1.3),
        Inches(3.35),
        Inches(2.3),
        "Pilotez votre établissement depuis un seul endroit",
        size=22,
        bold=True,
        color=WHITE,
    )
    add_text(
        s,
        panel + Inches(0.28),
        Inches(3.7),
        Inches(3.35),
        Inches(1.6),
        "Élèves, enseignants, présences, notes, finances et communication — sur Web et Mobile.",
        size=13,
        color=RGBColor(0xDB, 0xEA, 0xFE),
    )
    add_rect(s, panel + Inches(0.28), Inches(5.45), Inches(2.4), Inches(0.42), WHITE, rounded=True, radius=0.2)
    add_text(
        s,
        panel + Inches(0.28),
        Inches(5.47),
        Inches(2.4),
        Inches(0.38),
        "somafrik.app",
        size=13,
        bold=True,
        color=BRAND,
        align=PP_ALIGN.CENTER,
        anchor="m",
    )
    cover_shot = framed("leaf-cover-dash", WEB["dashboard"], max_w=900, radius=14)
    add_pic(s, cover_shot, panel + Inches(0.22), Inches(6.0), Inches(3.45))
    add_text(
        s,
        panel + Inches(0.22),
        Inches(7.85),
        Inches(3.45),
        Inches(0.28),
        "Transformer l’éducation, renforcer l’impact",
        size=10,
        color=RGBColor(0xDB, 0xEA, 0xFE),
        align=PP_ALIGN.CENTER,
    )

    # Rabat (droite)
    add_text(s, panel * 2 + Inches(0.28), Inches(0.35), Inches(3.35), Inches(0.28), "POURQUOI SOMAFRIK", size=11, bold=True, color=BRAND)
    add_text(
        s,
        panel * 2 + Inches(0.28),
        Inches(0.65),
        Inches(3.35),
        Inches(0.7),
        "Cinq apports concrets pour l’école",
        size=16,
        bold=True,
        color=INK,
    )
    for i, b in enumerate(BENEFITS):
        y = Inches(1.5) + i * Inches(1.05)
        add_rect(s, panel * 2 + Inches(0.28), y, Inches(3.35), Inches(0.92), WHITE, rounded=True, radius=0.12)
        add_rect(s, panel * 2 + Inches(0.28), y, Inches(0.1), Inches(0.92), BRAND)
        add_text(
            s,
            panel * 2 + Inches(0.5),
            y + Inches(0.18),
            Inches(3.0),
            Inches(0.6),
            b,
            size=12,
            color=INK,
            anchor="m",
        )
    add_text(
        s,
        panel * 2 + Inches(0.28),
        Inches(7.45),
        Inches(3.35),
        Inches(0.4),
        "Ouvrez le dépliant →",
        size=12,
        bold=True,
        color=BRAND,
    )

    fold_guides(s)

    # VERSO — intérieur : web | modules | mobile
    s = blank(prs)
    add_rect(s, 0, 0, W, H, WHITE)
    add_rect(s, 0, 0, W, Inches(0.08), BRAND)

    add_text(s, Inches(0.22), Inches(0.22), Inches(3.5), Inches(0.22), "APPLICATION WEB", size=10, bold=True, color=BRAND)
    add_text(s, Inches(0.22), Inches(0.44), Inches(3.55), Inches(0.7), "Le bureau pour piloter", size=16, bold=True, color=INK)
    add_text(
        s,
        Inches(0.22),
        Inches(1.1),
        Inches(3.5),
        Inches(0.55),
        "Tableaux de bord, classes et annuaire des élèves sur un écran complet.",
        size=11,
        color=MUTED,
    )
    shot = framed("leaf-dash", WEB["dashboard"], max_w=1000, radius=16)
    add_pic(s, shot, Inches(0.18), Inches(1.7), Inches(3.55))
    shot = framed("leaf-cls", WEB["classes"], max_w=1000, radius=16)
    add_pic(s, shot, Inches(0.18), Inches(4.55), Inches(3.55))

    add_text(
        s,
        panel + Inches(0.22),
        Inches(0.22),
        Inches(3.5),
        Inches(0.22),
        "AUJOURD’HUI DANS SOMAFRIK",
        size=10,
        bold=True,
        color=BRAND,
    )
    add_text(
        s,
        panel + Inches(0.22),
        Inches(0.44),
        Inches(3.55),
        Inches(0.55),
        "Les besoins couverts",
        size=16,
        bold=True,
        color=INK,
    )
    for i, (title, desc) in enumerate(FEATURES):
        y = Inches(1.1) + i * Inches(0.85)
        add_rect(s, panel + Inches(0.2), y, Inches(3.5), Inches(0.76), CANVAS, rounded=True, radius=0.12)
        add_text(s, panel + Inches(0.38), y + Inches(0.08), Inches(3.2), Inches(0.26), title, size=12, bold=True, color=INK)
        add_text(s, panel + Inches(0.38), y + Inches(0.34), Inches(3.2), Inches(0.36), desc, size=10, color=MUTED)

    add_text(
        s,
        panel + Inches(0.22),
        Inches(5.45),
        Inches(3.5),
        Inches(0.25),
        "POUR QUI ?",
        size=10,
        bold=True,
        color=BRAND,
    )
    add_text(
        s,
        panel + Inches(0.22),
        Inches(5.7),
        Inches(3.5),
        Inches(2.2),
        "Direction  ·  Administration\nEnseignants  ·  Parents\n\nChaque profil retrouve son travail quotidien, sans changer d’outil.",
        size=12,
        color=INK,
    )

    add_text(
        s,
        panel * 2 + Inches(0.18),
        Inches(0.22),
        Inches(3.5),
        Inches(0.22),
        "APPLICATION MOBILE",
        size=10,
        bold=True,
        color=TEAL,
    )
    add_text(
        s,
        panel * 2 + Inches(0.18),
        Inches(0.44),
        Inches(3.55),
        Inches(0.55),
        "Le terrain, dans la poche",
        size=16,
        bold=True,
        color=INK,
    )

    phones = [
        (MOBILE["classes"], "Classes"),
        (MOBILE["appel"], "Appel"),
        (MOBILE["notes"], "Notes"),
    ]
    pw = Inches(1.12)
    for i, (src, cap) in enumerate(phones):
        x = panel * 2 + Inches(0.16) + i * (pw + Inches(0.08))
        shot = framed(f"leaf-m-{i}", src, max_w=520, radius=28)
        add_pic(s, shot, x, Inches(1.15), pw)
        add_text(s, x, Inches(6.95), pw, Inches(0.28), cap, size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)

    add_text(
        s,
        panel * 2 + Inches(0.18),
        Inches(7.3),
        Inches(3.5),
        Inches(0.55),
        "Également : paiements, enseignants, espace parent.",
        size=10,
        color=MUTED,
    )

    fold_guides(s)

    out = DIST / "Somafrik-depliant-3-volets-Canva.pptx"
    prs.save(out)
    return out


# ---------------------------------------------------------------------------
# Aperçus PNG (contrôle visuel + fonds Canva)
# ---------------------------------------------------------------------------

def _t(draw, xy, text, fnt, fill, *, max_w=None):
    if max_w is None:
        draw.text(xy, text, font=fnt, fill=fill)
        return
    words = text.split()
    lines, cur = [], ""
    for w in words:
        trial = f"{cur} {w}".strip()
        if draw.textlength(trial, font=fnt) <= max_w:
            cur = trial
        else:
            if cur:
                lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    x, y = xy
    lh = int(fnt.size * 1.25)
    for line in lines:
        draw.text((x, y), line, font=fnt, fill=fill)
        y += lh


def paste_shot(base: Image.Image, src: Path, box: tuple[int, int, int, int], *, radius: int = 18) -> None:
    im = Image.open(src).convert("RGBA")
    tw, th = box[2] - box[0], box[3] - box[1]
    im.thumbnail((tw, th), Image.Resampling.LANCZOS)
    x = box[0] + (tw - im.width) // 2
    y = box[1]
    r = min(radius, im.width // 10, im.height // 10)
    mask = rounded_mask(im.size, r)
    im.putalpha(mask)
    base.alpha_composite(im, (x, y))


def render_presentation_previews() -> list[Path]:
    PREVIEWS.mkdir(parents=True, exist_ok=True)
    W, H = 1920, 1080
    out_paths: list[Path] = []

    # 1 cover
    im = Image.new("RGBA", (W, H), (*PIL_CANVAS, 255))
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, 820, H), fill=PIL_NAVY)
    d.rectangle((800, 0, 832, H), fill=PIL_BRAND)
    icon = Image.open(ICON).convert("RGBA").resize((88, 88), Image.Resampling.LANCZOS)
    im.alpha_composite(icon, (70, 55))
    _t(d, (175, 62), "SOMAFRIK", font(FONT_BOLD, 22), PIL_WHITE)
    _t(d, (175, 96), "Transformer l’éducation, renforcer l’impact", font(FONT_REG, 15), (203, 213, 225))
    _t(d, (70, 210), "GESTION SCOLAIRE MODERNE", font(FONT_BOLD, 16), (147, 197, 253))
    _t(d, (70, 260), "Pilotez votre établissement depuis un seul endroit", font(FONT_BOLD, 48), PIL_WHITE, max_w=700)
    _t(
        d,
        (70, 620),
        "Somafrik réunit les élèves, les enseignants, les présences, les notes, les finances et la communication dans une plateforme accessible sur Web et Mobile.",
        font(FONT_REG, 20),
        (203, 213, 225),
        max_w=700,
    )
    d.rounded_rectangle((70, 860, 430, 930), radius=14, fill=PIL_BRAND)
    _t(d, (145, 878), "somafrik.app", font(FONT_BOLD, 22), PIL_WHITE)
    _t(d, (460, 880), "Web  ·  Mobile", font(FONT_SEMI, 20), (147, 197, 253))
    paste_shot(im, WEB["dashboard"], (880, 110, 1860, 980), radius=16)
    p = PREVIEWS / "01-couverture.png"
    im.convert("RGB").save(p, "PNG")
    out_paths.append(p)

    # 2 web
    im = Image.new("RGBA", (W, H), (*PIL_WHITE, 255))
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, W, 12), fill=PIL_BRAND)
    _t(d, (70, 40), "APPLICATION WEB", font(FONT_BOLD, 16), PIL_BRAND)
    _t(d, (70, 78), "Ce que voit la direction une fois connectée", font(FONT_BOLD, 36), PIL_INK)
    _t(
        d,
        (70, 140),
        "Annuaire des élèves, organisation des classes, suivi des enseignants et tableaux de bord.",
        font(FONT_REG, 18),
        PIL_MUTED,
        max_w=1700,
    )
    cards = [(WEB["dashboard"], "Tableau de bord"), (WEB["classes"], "Classes"), (WEB["eleves"], "Annuaire des élèves")]
    for i, (src, cap) in enumerate(cards):
        x0 = 60 + i * 610
        d.rounded_rectangle((x0, 210, x0 + 580, 980), radius=18, fill=PIL_CANVAS)
        paste_shot(im, src, (x0 + 16, 226, x0 + 564, 900), radius=12)
        _t(d, (x0 + 180, 920), cap, font(FONT_SEMI, 20), PIL_INK)
    p = PREVIEWS / "02-application-web.png"
    im.convert("RGB").save(p, "PNG")
    out_paths.append(p)

    # 3 mobile
    im = Image.new("RGBA", (W, H), (*PIL_CANVAS, 255))
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, W, 12), fill=PIL_TEAL)
    _t(d, (70, 40), "APPLICATION MOBILE", font(FONT_BOLD, 16), PIL_TEAL)
    _t(d, (70, 78), "Le bureau pour piloter, le mobile pour le terrain", font(FONT_BOLD, 34), PIL_INK)
    _t(
        d,
        (70, 140),
        "Les enseignants font l’appel et saisissent les notes. Les parents suivent la scolarité de l’enfant.",
        font(FONT_REG, 18),
        PIL_MUTED,
        max_w=1700,
    )
    phones = [
        (MOBILE["classes"], "Classes"),
        (MOBILE["appel"], "Appel des présences"),
        (MOBILE["notes"], "Saisie des notes"),
        (MOBILE["parent"], "Espace parent"),
    ]
    for i, (src, cap) in enumerate(phones):
        x0 = 80 + i * 460
        d.rounded_rectangle((x0, 200, x0 + 420, 1000), radius=22, fill=PIL_WHITE)
        paste_shot(im, src, (x0 + 18, 218, x0 + 402, 940), radius=28)
        _t(d, (x0 + 110, 955), cap, font(FONT_SEMI, 18), PIL_INK)
    p = PREVIEWS / "03-application-mobile.png"
    im.convert("RGB").save(p, "PNG")
    out_paths.append(p)

    # 4 preuves
    im = Image.new("RGBA", (W, H), (*PIL_WHITE, 255))
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, W, 12), fill=PIL_BRAND)
    _t(d, (70, 40), "PREUVES MÉTIER", font(FONT_BOLD, 16), PIL_BRAND)
    _t(d, (70, 78), "Somafrik dans le quotidien de l’établissement", font(FONT_BOLD, 34), PIL_INK)
    proofs = [
        (WEB["login"], "Connexion web", "Code établissement et identifiant."),
        (WEB["dossier"], "Dossier élève", "Identité, inscription, responsables."),
        (MOBILE["paiements"], "Paiements", "Frais, encaissements, reçus."),
        (MOBILE["eleves"], "Élèves (mobile)", "Annuaire et suivi de classe."),
    ]
    for i, (src, title, desc) in enumerate(proofs):
        x0 = 50 + i * 465
        d.rounded_rectangle((x0, 170, x0 + 445, 1020), radius=18, fill=PIL_CANVAS)
        paste_shot(im, src, (x0 + 16, 186, x0 + 429, 860), radius=14)
        _t(d, (x0 + 24, 880), title, font(FONT_BOLD, 20), PIL_INK)
        _t(d, (x0 + 24, 920), desc, font(FONT_REG, 16), PIL_MUTED, max_w=400)
    p = PREVIEWS / "04-preuves-metier.png"
    im.convert("RGB").save(p, "PNG")
    out_paths.append(p)

    # 5 features
    im = Image.new("RGBA", (W, H), (*PIL_CANVAS, 255))
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, W, 12), fill=PIL_BRAND)
    _t(d, (70, 36), "FONCTIONNALITÉS", font(FONT_BOLD, 16), PIL_BRAND)
    _t(d, (70, 72), "Les besoins couverts aujourd’hui", font(FONT_BOLD, 34), PIL_INK)
    for i, (title, desc) in enumerate(FEATURES):
        x0 = 60 + i * 370
        d.rounded_rectangle((x0, 160, x0 + 350, 430), radius=16, fill=PIL_WHITE)
        d.rectangle((x0, 160, x0 + 10, 430), fill=PIL_BRAND)
        _t(d, (x0 + 28, 185), title, font(FONT_BOLD, 22), PIL_INK)
        _t(d, (x0 + 28, 235), desc, font(FONT_REG, 16), PIL_MUTED, max_w=300)
    _t(d, (70, 480), "POUR QUI ?", font(FONT_BOLD, 16), PIL_BRAND)
    _t(d, (70, 516), "Pensé d’abord pour l’établissement", font(FONT_BOLD, 26), PIL_INK)
    for i, (title, desc) in enumerate(AUDIENCES):
        x0 = 60 + i * 460
        d.rounded_rectangle((x0, 580, x0 + 440, 980), radius=16, fill=PIL_WHITE)
        _t(d, (x0 + 28, 610), title, font(FONT_BOLD, 24), PIL_BRAND)
        _t(d, (x0 + 28, 670), desc, font(FONT_REG, 18), PIL_INK, max_w=380)
    p = PREVIEWS / "05-fonctionnalites.png"
    im.convert("RGB").save(p, "PNG")
    out_paths.append(p)

    # 6 cta
    im = Image.new("RGBA", (W, H), (*PIL_NAVY, 255))
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, 16, H), fill=PIL_BRAND)
    icon = Image.open(ICON).convert("RGBA").resize((80, 80), Image.Resampling.LANCZOS)
    im.alpha_composite(icon, (80, 70))
    _t(d, (180, 80), "SOMAFRIK", font(FONT_BOLD, 22), PIL_WHITE)
    _t(d, (180, 116), "La plateforme qui simplifie la gestion de votre établissement scolaire.", font(FONT_REG, 16), (148, 163, 184))
    _t(d, (80, 230), "Ouvrez l’espace de votre établissement", font(FONT_BOLD, 42), PIL_WHITE, max_w=1050)
    _t(
        d,
        (80, 430),
        "Si votre école dispose déjà d’un compte Somafrik, connectez-vous pour retrouver les classes, les élèves, les présences, les notes et les finances.",
        font(FONT_REG, 20),
        (203, 213, 225),
        max_w=1000,
    )
    d.rounded_rectangle((80, 620, 520, 700), radius=14, fill=PIL_BRAND)
    _t(d, (130, 642), "somafrik.app /connexion", font(FONT_BOLD, 20), PIL_WHITE)
    _t(d, (550, 642), "somafrik.app", font(FONT_BOLD, 22), (147, 197, 253))
    d.rounded_rectangle((1220, 80, 1860, 960), radius=20, fill=(18, 26, 44))
    _t(d, (1260, 110), "Chaque établissement reste", font(FONT_BOLD, 20), PIL_WHITE)
    _t(d, (1260, 140), "dans son périmètre", font(FONT_BOLD, 20), PIL_WHITE)
    secs = [
        ("Authentification", "Connexion avec le code de l’établissement et un identifiant."),
        ("Séparation", "Les données d’une école ne sont pas mélangées avec celles d’une autre."),
        ("Rôles", "Direction, administration, enseignant ou parent : chacun voit son périmètre."),
    ]
    for i, (t, desc) in enumerate(secs):
        y = 220 + i * 200
        _t(d, (1260, y), t, font(FONT_BOLD, 20), (147, 197, 253))
        _t(d, (1260, y + 40), desc, font(FONT_REG, 17), (203, 213, 225), max_w=540)
    _t(d, (80, 1000), "© 2026 Somafrik  ·  Captures runtime, données fictives  ·  Document modifiable dans Canva", font(FONT_REG, 14), (100, 116, 139))
    p = PREVIEWS / "06-appel-action.png"
    im.convert("RGB").save(p, "PNG")
    out_paths.append(p)
    return out_paths


def render_leaflet_previews() -> list[Path]:
    PREVIEWS.mkdir(parents=True, exist_ok=True)
    W, H = 3508, 2480  # A4 paysage ~300 dpi
    pw = W // 3
    paths: list[Path] = []

    im = Image.new("RGBA", (W, H), (*PIL_WHITE, 255))
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, pw, H), fill=PIL_NAVY)
    d.rectangle((pw, 0, pw * 2, H), fill=PIL_BRAND)
    d.rectangle((pw * 2, 0, W, H), fill=PIL_CANVAS)
    icon = Image.open(ICON).convert("RGBA").resize((110, 110), Image.Resampling.LANCZOS)
    im.alpha_composite(icon, (70, 80))
    _t(d, (200, 95), "SOMAFRIK", font(FONT_BOLD, 36), PIL_WHITE)
    _t(d, (200, 145), "Web et mobile", font(FONT_MED, 22), (147, 197, 253))
    _t(d, (70, 280), "Contact", font(FONT_BOLD, 48), PIL_WHITE)
    _t(d, (70, 380), "somafrik.app", font(FONT_SEMI, 32), PIL_WHITE)
    _t(d, (70, 430), "somafrik.app/connexion", font(FONT_REG, 26), (226, 232, 240))
    _t(d, (70, 560), "Sécurité", font(FONT_BOLD, 32), (147, 197, 253))
    _t(d, (70, 630), "Authentification par établissement.", font(FONT_REG, 26), (203, 213, 225), max_w=pw - 140)
    _t(d, (70, 760), "Données séparées par école.", font(FONT_REG, 26), (203, 213, 225), max_w=pw - 140)
    _t(d, (70, 890), "Accès selon le rôle : direction, administration, enseignant, parent.", font(FONT_REG, 26), (203, 213, 225), max_w=pw - 140)
    _t(d, (70, 2280), "© 2026 Somafrik", font(FONT_REG, 22), (100, 116, 139))
    _t(d, (70, 2320), "Captures runtime · données fictives", font(FONT_REG, 20), (100, 116, 139))

    im.alpha_composite(icon, (pw + 70, 80))
    _t(d, (pw + 210, 95), "SOMAFRIK", font(FONT_BOLD, 36), PIL_WHITE)
    _t(d, (pw + 210, 145), "Gestion scolaire moderne", font(FONT_MED, 22), (219, 234, 254))
    _t(d, (pw + 70, 320), "Pilotez votre établissement depuis un seul endroit", font(FONT_BOLD, 58), PIL_WHITE, max_w=pw - 140)
    _t(
        d,
        (pw + 70, 900),
        "Élèves, enseignants, présences, notes, finances et communication — sur Web et Mobile.",
        font(FONT_REG, 28),
        (219, 234, 254),
        max_w=pw - 140,
    )
    d.rounded_rectangle((pw + 70, 1240, pw + 520, 1340), radius=20, fill=PIL_WHITE)
    _t(d, (pw + 160, 1265), "somafrik.app", font(FONT_BOLD, 32), PIL_BRAND)
    paste_shot(im, WEB["dashboard"], (pw + 60, 1380, pw * 2 - 60, 2140), radius=16)
    _t(d, (pw + 70, 2180), "Transformer l’éducation, renforcer l’impact", font(FONT_SEMI, 22), (219, 234, 254))

    _t(d, (pw * 2 + 70, 80), "POURQUOI SOMAFRIK", font(FONT_BOLD, 22), PIL_BRAND)
    _t(d, (pw * 2 + 70, 130), "Cinq apports concrets pour l’école", font(FONT_BOLD, 36), PIL_INK, max_w=pw - 140)
    for i, b in enumerate(BENEFITS):
        y = 380 + i * 280
        d.rounded_rectangle((pw * 2 + 60, y, W - 70, y + 240), radius=22, fill=PIL_WHITE)
        d.rectangle((pw * 2 + 60, y, pw * 2 + 82, y + 240), fill=PIL_BRAND)
        _t(d, (pw * 2 + 110, y + 70), b, font(FONT_SEMI, 26), PIL_INK, max_w=pw - 220)
    _t(d, (pw * 2 + 70, 2280), "Ouvrez le dépliant →", font(FONT_BOLD, 26), PIL_BRAND)
    for x in (pw, pw * 2):
        d.rectangle((x - 3, 0, x + 3, H), fill=(203, 213, 225))
    p = PREVIEWS / "depliant-recto-exterieur.png"
    im.convert("RGB").save(p, "PNG", optimize=True)
    paths.append(p)

    im = Image.new("RGBA", (W, H), (*PIL_WHITE, 255))
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, W, 18), fill=PIL_BRAND)
    _t(d, (50, 50), "APPLICATION WEB", font(FONT_BOLD, 22), PIL_BRAND)
    _t(d, (50, 95), "Le bureau pour piloter", font(FONT_BOLD, 40), PIL_INK)
    _t(d, (50, 170), "Tableaux de bord, classes et annuaire des élèves sur un écran complet.", font(FONT_REG, 24), PIL_MUTED, max_w=pw - 100)
    paste_shot(im, WEB["dashboard"], (40, 320, pw - 40, 1280), radius=16)
    paste_shot(im, WEB["classes"], (40, 1320, pw - 40, 2360), radius=16)

    _t(d, (pw + 50, 50), "AUJOURD’HUI DANS SOMAFRIK", font(FONT_BOLD, 22), PIL_BRAND)
    _t(d, (pw + 50, 95), "Les besoins couverts", font(FONT_BOLD, 40), PIL_INK)
    for i, (title, desc) in enumerate(FEATURES):
        y = 220 + i * 230
        d.rounded_rectangle((pw + 45, y, pw * 2 - 45, y + 200), radius=18, fill=PIL_CANVAS)
        _t(d, (pw + 80, y + 35), title, font(FONT_BOLD, 28), PIL_INK)
        _t(d, (pw + 80, y + 90), desc, font(FONT_REG, 22), PIL_MUTED, max_w=pw - 180)
    _t(d, (pw + 50, 1420), "POUR QUI ?", font(FONT_BOLD, 22), PIL_BRAND)
    _t(
        d,
        (pw + 50, 1480),
        "Direction  ·  Administration  ·  Enseignants  ·  Parents. Chaque profil retrouve son travail quotidien, sans changer d’outil.",
        font(FONT_REG, 26),
        PIL_INK,
        max_w=pw - 120,
    )

    _t(d, (pw * 2 + 40, 50), "APPLICATION MOBILE", font(FONT_BOLD, 22), PIL_TEAL)
    _t(d, (pw * 2 + 40, 95), "Le terrain, dans la poche", font(FONT_BOLD, 36), PIL_INK)
    phones = [(MOBILE["classes"], "Classes"), (MOBILE["appel"], "Appel"), (MOBILE["notes"], "Notes")]
    col_w = (pw - 80) // 3
    for i, (src, cap) in enumerate(phones):
        x0 = pw * 2 + 30 + i * (col_w + 10)
        paste_shot(im, src, (x0, 220, x0 + col_w, 2100), radius=28)
        _t(d, (x0 + 40, 2140), cap, font(FONT_SEMI, 24), PIL_INK)
    _t(d, (pw * 2 + 40, 2260), "Également : paiements, enseignants, espace parent.", font(FONT_REG, 22), PIL_MUTED, max_w=pw - 80)
    for x in (pw, pw * 2):
        d.rectangle((x - 3, 0, x + 3, H), fill=(203, 213, 225))
    p = PREVIEWS / "depliant-verso-interieur.png"
    im.convert("RGB").save(p, "PNG", optimize=True)
    paths.append(p)
    return paths


def build_pdf(preview_paths: list[Path], dest: Path) -> Path:
    rgb = [Image.open(p).convert("RGB") for p in preview_paths]
    rgb[0].save(dest, "PDF", save_all=True, append_images=rgb[1:], resolution=150)
    return dest


def main() -> None:
    Path("/tmp/somafrik-depliant").mkdir(parents=True, exist_ok=True)
    PREVIEWS.mkdir(parents=True, exist_ok=True)
    pres = build_presentation()
    leaflet = build_leaflet()
    slides = render_presentation_previews()
    folds = render_leaflet_previews()
    pdf_pres = build_pdf(slides, DIST / "Somafrik-presentation.pdf")
    pdf_leaf = build_pdf(folds, DIST / "Somafrik-depliant-3-volets.pdf")
    print("OK")
    for p in (pres, leaflet, pdf_pres, pdf_leaf, *slides, *folds):
        print(f"  {p}  ({p.stat().st_size // 1024}K)")


if __name__ == "__main__":
    main()
